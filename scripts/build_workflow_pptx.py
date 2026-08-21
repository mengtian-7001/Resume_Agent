#!/usr/bin/env python3
"""Export the dual-agent workflow as a 16:9 PPT matching the 简历中台 frontend."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from pptx.slide import Slide
from pptx.shapes.autoshape import Shape

INK = RGBColor(0x11, 0x11, 0x11)
MUTED = RGBColor(0x6E, 0x6E, 0x73)
LINE = RGBColor(0xE5, 0xE5, 0xE7)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
MINT = RGBColor(0xF5, 0xF5, 0xF7)
SOFT = RGBColor(0xFA, 0xFA, 0xFB)
HAIR = RGBColor(0xD2, 0xD2, 0xD7)
FONT = "PingFang SC"


def emu(value) -> int:
    return int(round(float(value)))


def _set_typeface(run, name: str = FONT) -> None:
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def style_run(run, *, size: int, bold: bool = False, color: RGBColor = INK) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    _set_typeface(run)


def fill_rgb(shape: Shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def line_rgb(shape: Shape, color: RGBColor, pt: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(pt)


def no_line(shape: Shape) -> None:
    shape.line.fill.background()


def add_round(
    slide: Slide,
    x,
    y,
    w,
    h,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    weight: float = 1.0,
    adj: float = 0.12,
) -> Shape:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    fill_rgb(shape, fill)
    if line is None:
        no_line(shape)
    else:
        line_rgb(shape, line, weight)
    try:
        shape.adjustments[0] = adj
    except Exception:
        pass
    return shape


def write_shape(
    shape: Shape,
    lines: list[tuple[str, int, bool, RGBColor]],
    *,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.MIDDLE,
    pad: float = 0.08,
) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(pad)
    tf.margin_right = Inches(pad)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.TOP: "t"}.get(anchor, "ctr"))
    except Exception:
        pass
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(1 if i < len(lines) - 1 else 0)
        p.clear()
        run = p.add_run()
        run.text = text
        style_run(run, size=size, bold=bold, color=color)


def add_label(slide: Slide, x, y, w, h, text: str, *, size=11, bold=False, color=MUTED, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE) -> Shape:
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    write_shape(box, [(text, size, bold, color)], align=align, anchor=anchor, pad=0.02)
    return box


def add_chip(slide: Slide, x, y, w, h, title: str, detail: str, *, accent: bool = False) -> Shape:
    shape = add_round(
        slide,
        x,
        y,
        w,
        h,
        fill=INK if accent else PAPER,
        line=INK if accent else LINE,
        weight=1.0,
        adj=0.16,
    )
    write_shape(
        shape,
        [
            (title, 11, True, PAPER if accent else INK),
            (detail, 9, False, RGBColor(0xC7, 0xC7, 0xCC) if accent else MUTED),
        ],
    )
    return shape


def add_step(slide: Slide, x, y, w, h, n: str, title: str) -> Shape:
    shape = add_round(slide, x, y, w, h, fill=PAPER, line=LINE, adj=0.14)
    write_shape(
        shape,
        [(n, 8, False, MUTED), (title, 10, True, INK)],
    )
    return shape


def down_arrow(slide: Slide, cx, y, label: str) -> None:
    add_label(slide, cx - Inches(1.6), y, Inches(3.2), Inches(0.18), label, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    stem = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        emu(cx - Inches(0.01)),
        emu(y + Inches(0.18)),
        emu(Inches(0.02)),
        emu(Inches(0.10)),
    )
    fill_rgb(stem, INK)
    no_line(stem)
    head = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        emu(cx - Inches(0.07)),
        emu(y + Inches(0.26)),
        emu(Inches(0.14)),
        emu(Inches(0.10)),
    )
    fill_rgb(head, INK)
    no_line(head)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    cover(prs, blank)
    flow(prs, blank)
    harness(prs, blank)
    scoring(prs, blank)
    scoring_sources(prs, blank)
    return prs


def cover(prs: Presentation, blank) -> None:
    slide = prs.slides.add_slide(blank)
    fill_rgb(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, emu(prs.slide_width), emu(prs.slide_height)), PAPER)
    no_line(slide.shapes[-1])

    mark = add_round(slide, Inches(0.72), Inches(0.55), Inches(0.38), Inches(0.38), fill=INK, adj=0.5)
    write_shape(mark, [("简", 14, True, PAPER)], align=PP_ALIGN.CENTER, pad=0.0)
    add_label(slide, Inches(1.22), Inches(0.55), Inches(5), Inches(0.22), "简历中台", size=14, bold=True, color=INK, anchor=MSO_ANCHOR.TOP)
    add_label(slide, Inches(1.22), Inches(0.76), Inches(6), Inches(0.2), "智能筛选工作台", size=11, color=MUTED, anchor=MSO_ANCHOR.TOP)

    add_label(
        slide,
        Inches(0.72),
        Inches(2.55),
        Inches(11.5),
        Inches(0.9),
        "双 Agent 筛选工作流",
        size=40,
        bold=True,
        color=INK,
        anchor=MSO_ANCHOR.TOP,
    )
    add_label(
        slide,
        Inches(0.72),
        Inches(3.5),
        Inches(10.2),
        Inches(0.7),
        "模型读懂文本，Harness 管住边界，两个 Agent 一产一审，人做最终决定。",
        size=16,
        color=MUTED,
        anchor=MSO_ANCHOR.TOP,
    )

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        emu(Inches(0.72)),
        emu(Inches(6.55)),
        emu(Inches(11.9)),
        emu(Inches(0.015)),
    )
    fill_rgb(rule, LINE)
    no_line(rule)
    add_label(
        slide,
        Inches(0.72),
        Inches(6.7),
        Inches(11.9),
        Inches(0.3),
        "系统可以错，但不能装对。录用权在人。",
        size=13,
        color=INK,
        anchor=MSO_ANCHOR.TOP,
    )


def flow(prs: Presentation, blank) -> None:
    slide = prs.slides.add_slide(blank)
    fill_rgb(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, emu(prs.slide_width), emu(prs.slide_height)), PAPER)
    no_line(slide.shapes[-1])

    mark = add_round(slide, Inches(0.42), Inches(0.22), Inches(0.28), Inches(0.28), fill=INK, adj=0.5)
    write_shape(mark, [("简", 11, True, PAPER)], align=PP_ALIGN.CENTER, pad=0.0)
    add_label(slide, Inches(0.78), Inches(0.18), Inches(4.2), Inches(0.18), "简历中台 · 筛选工作流", size=13, bold=True, color=INK, anchor=MSO_ANCHOR.TOP)
    add_label(
        slide,
        Inches(0.78),
        Inches(0.36),
        Inches(7.5),
        Inches(0.18),
        "模型读懂文本 · Harness 管住边界 · 两个 Agent 一产一审 · 人做最终决定",
        size=10,
        color=MUTED,
        anchor=MSO_ANCHOR.TOP,
    )

    pill = add_round(slide, Inches(9.55), Inches(0.24), Inches(1.35), Inches(0.28), fill=INK, adj=0.5)
    write_shape(pill, [("黑框 = Agent", 9, True, PAPER)], align=PP_ALIGN.CENTER, pad=0.04)
    pill2 = add_round(slide, Inches(11.0), Inches(0.24), Inches(1.9), Inches(0.28), fill=MINT, line=LINE, adj=0.5)
    write_shape(pill2, [("灰框 = Harness / 规则", 9, False, MUTED)], align=PP_ALIGN.CENTER, pad=0.04)

    hair = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        emu(Inches(0.42)),
        emu(Inches(0.64)),
        emu(Inches(12.5)),
        emu(Inches(0.012)),
    )
    fill_rgb(hair, LINE)
    no_line(hair)

    lx = Inches(0.42)
    cx = Inches(1.18)
    cw = Inches(11.72)
    gap = Inches(0.08)
    chip_h = Inches(0.58)

    def lane_shell(y, h, name: str) -> None:
        add_label(slide, lx, y, Inches(0.7), h, name, size=11, color=MUTED, align=PP_ALIGN.LEFT)
        add_round(slide, cx, y, cw, h, fill=SOFT, line=LINE, weight=1.0, adj=0.04)

    # 接入
    y = Inches(0.78)
    h = Inches(0.98)
    lane_shell(y, h, "接入")
    chip_w = (cw - Inches(0.28) - gap * 2) / 3
    row_y = y + Inches(0.1)
    x0 = cx + Inches(0.14)
    add_chip(slide, x0, row_y, chip_w, chip_h, "招聘人员上传", "1 份 JD + N 份简历")
    add_chip(slide, x0 + chip_w + gap, row_y, chip_w, chip_h, "入口校验", "格式 · 份数 · 大小")
    add_chip(slide, x0 + (chip_w + gap) * 2, row_y, chip_w, chip_h, "当场拒收", "不合法，不进模型")
    add_label(slide, x0, y + Inches(0.72), cw, Inches(0.2), "坏文件直接拒。过关后才开始理解内容。", size=9, color=MUTED)

    mid = cx + cw / 2
    down_arrow(slide, mid, y + h + Inches(0.02), "合法输入")

    # 准备
    y = Inches(2.18)
    h = Inches(1.18)
    lane_shell(y, h, "准备")
    row1 = y + Inches(0.1)
    add_chip(slide, x0, row1, chip_w, chip_h, "抽出岗位标准", "必备 / 年限 / 学历")
    add_chip(slide, x0 + chip_w + gap, row1, chip_w, chip_h, "补岗位背景", "术语语境，不证明候选人")
    add_chip(slide, x0 + (chip_w + gap) * 2, row1, chip_w, chip_h, "冻结 hard_gate", "全员同一把尺子")
    row2 = row1 + chip_h + Inches(0.07)
    half = (cw - Inches(0.28) - gap) / 2
    add_chip(slide, x0, row2, half, Inches(0.36), "Fan-out", "按简历并发，独立失败")
    add_chip(slide, x0 + half + gap, row2, half, Inches(0.36), "抽出候选人画像", "技能归一 · 字段带引用")

    down_arrow(slide, mid, y + h + Inches(0.02), "每个候选人进入双 Agent")

    # Agent
    y = Inches(3.80)
    h = Inches(2.48)
    lane_shell(y, h, "Agent")
    panel_y = y + Inches(0.1)
    panel_h = Inches(2.10)
    gutter = Inches(0.78)
    panel_w = (cw - Inches(0.28) - gutter) / 2
    left_x = x0
    right_x = x0 + panel_w + gutter

    def agent_panel(x, title: str, role: str) -> None:
        add_round(slide, x, panel_y, panel_w, panel_h, fill=MINT, line=INK, weight=1.5, adj=0.05)
        add_label(slide, x + Inches(0.12), panel_y + Inches(0.06), panel_w - Inches(0.24), Inches(0.22), title, size=14, bold=True, color=INK)
        add_label(slide, x + Inches(0.12), panel_y + Inches(0.26), panel_w - Inches(0.24), Inches(0.18), role, size=10, color=MUTED)

    agent_panel(left_x, "Construction 构建", "写稿：把这个人讲清楚")
    agent_panel(right_x, "Checker 审校", "质检：这句话站不站得住")

    step_y = panel_y + Inches(0.50)
    step_h = Inches(0.50)
    step_gap = Inches(0.05)
    step_w = (panel_w - Inches(0.24) - step_gap * 3) / 4
    for i, (n, title) in enumerate(
        [("1", "Plan 规划"), ("2", "Act 取证"), ("3", "Reflect"), ("4", "出题决策")]
    ):
        add_step(slide, left_x + Inches(0.12) + i * (step_w + step_gap), step_y, step_w, step_h, n, title)
    add_label(
        slide,
        left_x + Inches(0.12),
        step_y + step_h + Inches(0.04),
        panel_w - Inches(0.24),
        Inches(0.2),
        "内环：证据不足 → 再调白名单工具；预算用尽必须停。",
        size=9,
        color=MUTED,
    )
    mini_y = step_y + step_h + Inches(0.26)
    mini_w = (panel_w - Inches(0.24) - gap * 2) / 3
    mini_h = Inches(0.42)
    add_chip(slide, left_x + Inches(0.12), mini_y, mini_w, mini_h, "硬门槛", "模型抬不过")
    add_chip(slide, left_x + Inches(0.12) + mini_w + gap, mini_y, mini_w, mini_h, "规则分", "可解释锚点")
    add_chip(slide, left_x + Inches(0.12) + (mini_w + gap) * 2, mini_y, mini_w, mini_h, "模型评委", "必须引用原文")
    add_label(
        slide,
        left_x + Inches(0.12),
        panel_y + panel_h - Inches(0.26),
        panel_w - Inches(0.24),
        Inches(0.2),
        "产出：优先级 · 原文证据 · ≥10 题 + 追问",
        size=9,
        color=MUTED,
    )

    chk_w = (panel_w - Inches(0.24) - step_gap * 2) / 3
    for i, (n, title) in enumerate(
        [("A", "证据能否定位"), ("B", "分与结论一致"), ("C", "题目能否验证")]
    ):
        add_step(slide, right_x + Inches(0.12) + i * (chk_w + step_gap), step_y, chk_w, step_h, n, title)
    add_label(
        slide,
        right_x + Inches(0.12),
        step_y + step_h + Inches(0.04),
        panel_w - Inches(0.24),
        Inches(0.2),
        "拦住「了解→精通」、缺证据的推荐、不到 10 道题。",
        size=9,
        color=MUTED,
    )
    chk_mini_w = (panel_w - Inches(0.24) - gap) / 2
    add_chip(slide, right_x + Inches(0.12), mini_y, chk_mini_w, mini_h, "本地 patch", "降级 / 加风险 / 补题")
    add_chip(slide, right_x + Inches(0.12) + chk_mini_w + gap, mini_y, chk_mini_w, mini_h, "fail-closed", "质检挂了就复核")
    add_label(
        slide,
        right_x + Inches(0.12),
        panel_y + panel_h - Inches(0.26),
        panel_w - Inches(0.24),
        Inches(0.2),
        "产出：问题清单 · 严重度 · 安全补丁",
        size=9,
        color=MUTED,
    )

    # middle handshake
    mx = left_x + panel_w
    add_label(slide, mx, panel_y + Inches(0.72), gutter, Inches(0.22), "交卷 →", size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)
    dash = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        emu(mx + Inches(0.12)),
        emu(panel_y + Inches(1.08)),
        emu(mx + gutter - Inches(0.12)),
        emu(panel_y + Inches(1.08)),
    )
    dash.line.color.rgb = INK
    dash.line.width = Pt(1.15)
    dash.line.dash_style = MSO_LINE.DASH
    add_label(slide, mx, panel_y + Inches(1.16), gutter, Inches(0.22), "← 回修 1 次", size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)

    add_label(
        slide,
        x0,
        y + h - Inches(0.24),
        cw,
        Inches(0.2),
        "共享合同和原文，不共享决策权。虚线最多回一次，不继续改到「看起来通过」。",
        size=9,
        color=MUTED,
    )

    down_arrow(slide, mid, y + h + Inches(0.0), "通过，或已达回修上限")

    # 收口
    y = Inches(6.68)
    h = Inches(0.64)
    lane_shell(y, h, "收口")
    end_h = Inches(0.44)
    end_y = y + Inches(0.1)
    add_chip(slide, x0, end_y, chip_w, end_h, "Policy 只降不升", "recommend → review")
    add_chip(slide, x0 + chip_w + gap, end_y, chip_w, end_h, "交到工作台", "分数 · 证据 · 轨迹")
    add_chip(slide, x0 + (chip_w + gap) * 2, end_y, chip_w, end_h, "人做最终决定", "验证自述 · 团队适配", accent=True)


def harness(prs: Presentation, blank) -> None:
    slide = prs.slides.add_slide(blank)
    fill_rgb(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, emu(prs.slide_width), emu(prs.slide_height)), PAPER)
    no_line(slide.shapes[-1])

    mark = add_round(slide, Inches(0.42), Inches(0.22), Inches(0.28), Inches(0.28), fill=INK, adj=0.5)
    write_shape(mark, [("简", 11, True, PAPER)], align=PP_ALIGN.CENTER, pad=0.0)
    add_label(slide, Inches(0.78), Inches(0.18), Inches(6.2), Inches(0.18), "简历中台 · Harness 控制流", size=13, bold=True, color=INK, anchor=MSO_ANCHOR.TOP)
    add_label(
        slide,
        Inches(0.78),
        Inches(0.36),
        Inches(8.8),
        Inches(0.18),
        "模型在哪些地方不准自由发挥：不让进 · 不让乱写 · 不让无限跑 · 不让装对",
        size=10,
        color=MUTED,
        anchor=MSO_ANCHOR.TOP,
    )
    pill = add_round(slide, Inches(10.55), Inches(0.24), Inches(2.35), Inches(0.28), fill=INK, adj=0.5)
    write_shape(pill, [("强调框 = 运行时外壳", 9, True, PAPER)], align=PP_ALIGN.CENTER, pad=0.04)

    hair = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        emu(Inches(0.42)),
        emu(Inches(0.64)),
        emu(Inches(12.5)),
        emu(Inches(0.012)),
    )
    fill_rgb(hair, LINE)
    no_line(hair)

    lx = Inches(0.42)
    cx = Inches(1.18)
    cw = Inches(11.72)
    gap = Inches(0.08)
    chip_h = Inches(0.52)
    x0 = cx + Inches(0.14)
    chip_w = (cw - Inches(0.28) - gap * 2) / 3
    mid = cx + cw / 2

    def lane_shell(y, h, name: str) -> None:
        add_label(slide, lx, y, Inches(0.7), h, name, size=11, color=MUTED, align=PP_ALIGN.LEFT)
        add_round(slide, cx, y, cw, h, fill=SOFT, line=LINE, weight=1.0, adj=0.04)

    y = Inches(0.78)
    h = Inches(0.92)
    lane_shell(y, h, "进门")
    add_chip(slide, x0, y + Inches(0.1), chip_w, chip_h, "入口校验", "格式 · 份数 · 大小")
    add_chip(slide, x0 + chip_w + gap, y + Inches(0.1), chip_w, chip_h, "隔离不可信文本", "JD/简历只当 DATA")
    add_chip(slide, x0 + (chip_w + gap) * 2, y + Inches(0.1), chip_w, chip_h, "当场拒收", "不合法，不调模型")
    add_label(slide, x0, y + Inches(0.66), cw, Inches(0.2), "坏文件直接拒。模型还没出场。", size=9, color=MUTED)
    down_arrow(slide, mid, y + h + Inches(0.0), "合法输入才建合同")

    y = Inches(2.10)
    h = Inches(0.92)
    lane_shell(y, h, "契约")
    add_chip(slide, x0, y + Inches(0.1), chip_w, chip_h, "冻结岗位标准", "hard_gate 不可改写")
    add_chip(slide, x0 + chip_w + gap, y + Inches(0.1), chip_w, chip_h, "固定 Schema", "Construction / Checker 合同")
    add_chip(slide, x0 + (chip_w + gap) * 2, y + Inches(0.1), chip_w, chip_h, "只写自己的字段", "不能互相覆盖决策")
    add_label(slide, x0, y + Inches(0.66), cw, Inches(0.2), "先把能写什么、不能改什么写死。后面的智能只能在合同里填空。", size=9, color=MUTED)
    down_arrow(slide, mid, y + h + Inches(0.0), "带着预算和白名单进入 Agent")

    y = Inches(3.42)
    h = Inches(2.28)
    lane_shell(y, h, "运行时")
    shell = add_round(
        slide,
        x0,
        y + Inches(0.1),
        cw - Inches(0.28),
        Inches(2.06),
        fill=MINT,
        line=INK,
        weight=1.5,
        adj=0.04,
    )
    shell.text_frame.paragraphs[0].text = ""
    add_label(
        slide,
        x0 + Inches(0.12),
        y + Inches(0.12),
        cw - Inches(0.5),
        Inches(0.22),
        "外壳包住两个 Agent，不让它们裸跑",
        size=13,
        bold=True,
        color=INK,
    )
    ctrl_y = y + Inches(0.38)
    ctrl_h = Inches(0.46)
    ctrl_w = (cw - Inches(0.52) - gap * 2) / 3
    add_chip(slide, x0 + Inches(0.12), ctrl_y, ctrl_w, ctrl_h, "Budget", "步数 · LLM 次数 · 截止 · 租约", accent=True)
    add_chip(slide, x0 + Inches(0.12) + ctrl_w + gap, ctrl_y, ctrl_w, ctrl_h, "Policy", "7 个白名单工具 · 网页不当证据", accent=True)
    add_chip(slide, x0 + Inches(0.12) + (ctrl_w + gap) * 2, ctrl_y, ctrl_w, ctrl_h, "Observe", "agent_runs 逐步轨迹", accent=True)

    box_y = y + Inches(0.92)
    box_h = Inches(0.78)
    gutter = Inches(0.78)
    box_w = (cw - Inches(0.52) - gutter) / 2
    left_x = x0 + Inches(0.12)
    right_x = left_x + box_w + gutter
    left = add_round(slide, left_x, box_y, box_w, box_h, fill=PAPER, line=LINE, adj=0.08)
    write_shape(
        left,
        [("Construction", 12, True, INK), ("规划、取证、出题。工具失败就跳过，用规则分兜底。", 9, False, MUTED)],
        anchor=MSO_ANCHOR.TOP,
        pad=0.1,
    )
    right = add_round(slide, right_x, box_y, box_w, box_h, fill=PAPER, line=LINE, adj=0.08)
    write_shape(
        right,
        [("Checker", 12, True, INK), ("对照原文挑刺。不能发明经历，不能把红灯改绿灯。", 9, False, MUTED)],
        anchor=MSO_ANCHOR.TOP,
        pad=0.1,
    )
    mx = left_x + box_w
    add_label(slide, mx, box_y + Inches(0.12), gutter, Inches(0.2), "交卷 →", size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)
    dash = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        emu(mx + Inches(0.12)),
        emu(box_y + Inches(0.40)),
        emu(mx + gutter - Inches(0.12)),
        emu(box_y + Inches(0.40)),
    )
    dash.line.color.rgb = INK
    dash.line.width = Pt(1.15)
    dash.line.dash_style = MSO_LINE.DASH
    add_label(slide, mx, box_y + Inches(0.46), gutter, Inches(0.2), "← 最多 1 次", size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_label(
        slide,
        x0 + Inches(0.12),
        y + Inches(1.78),
        cw - Inches(0.5),
        Inches(0.28),
        "回环有上限：证据不足可以再找，质检不通过最多回修一次。预算用尽必须停。",
        size=9,
        color=MUTED,
    )
    down_arrow(slide, mid, y + h + Inches(0.0), "交卷后先验证，再决定能不能见人")

    y = Inches(6.12)
    h = Inches(1.16)
    lane_shell(y, h, "出门")
    four_w = (cw - Inches(0.28) - gap * 3) / 4
    four_h = Inches(0.58)
    fy = y + Inches(0.1)
    add_chip(slide, x0, fy, four_w, four_h, "Verify", "引用定位 · 强度 · 题结构")
    add_chip(slide, x0 + four_w + gap, fy, four_w, four_h, "Degrade", "recommend 只降不升")
    add_chip(slide, x0 + (four_w + gap) * 2, fy, four_w, four_h, "fail-closed", "质检挂了就复核")
    add_chip(slide, x0 + (four_w + gap) * 3, fy, four_w, four_h, "交给人", "摊开证据和保留意见", accent=True)
    add_label(
        slide,
        x0,
        y + Inches(0.76),
        cw,
        Inches(0.3),
        "Harness 的最后一件事不是让结论更漂亮，而是不让系统装对。",
        size=9,
        color=MUTED,
    )


def scoring(prs: Presentation, blank) -> None:
    slide = prs.slides.add_slide(blank)
    fill_rgb(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, emu(prs.slide_width), emu(prs.slide_height)), PAPER)
    no_line(slide.shapes[-1])

    mark = add_round(slide, Inches(0.42), Inches(0.22), Inches(0.28), Inches(0.28), fill=INK, adj=0.5)
    write_shape(mark, [("简", 11, True, PAPER)], align=PP_ALIGN.CENTER, pad=0.0)
    add_label(slide, Inches(0.78), Inches(0.18), Inches(6.5), Inches(0.18), "简历中台 · 匹配打分怎么来的", size=13, bold=True, color=INK, anchor=MSO_ANCHOR.TOP)
    add_label(
        slide,
        Inches(0.78),
        Inches(0.36),
        Inches(9.2),
        Inches(0.18),
        "对齐开源 ATS：硬门槛单独拦，双路打分，0.60 / 0.40 合成。不是自拟权重。",
        size=10,
        color=MUTED,
        anchor=MSO_ANCHOR.TOP,
    )
    pill = add_round(slide, Inches(10.35), Inches(0.24), Inches(2.55), Inches(0.28), fill=INK, adj=0.5)
    write_shape(pill, [("强调框 = 开源公式落地", 9, True, PAPER)], align=PP_ALIGN.CENTER, pad=0.04)

    hair = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        emu(Inches(0.42)),
        emu(Inches(0.64)),
        emu(Inches(12.5)),
        emu(Inches(0.012)),
    )
    fill_rgb(hair, LINE)
    no_line(hair)

    lx = Inches(0.38)
    cx = Inches(1.22)
    cw = Inches(11.68)
    gap = Inches(0.08)
    x0 = cx + Inches(0.14)
    chip_w = (cw - Inches(0.28) - gap * 2) / 3
    mid = cx + cw / 2

    def lane_shell(y, h, name: str) -> None:
        add_label(slide, lx, y, Inches(0.78), h, name, size=11, color=MUTED, align=PP_ALIGN.LEFT)
        add_round(slide, cx, y, cw, h, fill=SOFT, line=LINE, weight=1.0, adj=0.04)

    y = Inches(0.76)
    h = Inches(0.86)
    lane_shell(y, h, "依据")
    add_chip(slide, x0, y + Inches(0.1), chip_w, Inches(0.50), "ema-resume-ranker", "0.60 LLM + 0.40 确定性", accent=True)
    add_chip(slide, x0 + chip_w + gap, y + Inches(0.1), chip_w, Inches(0.50), "ResuRank", "文本：0.60 语义 + 0.40 TF-IDF")
    add_chip(slide, x0 + (chip_w + gap) * 2, y + Inches(0.1), chip_w, Inches(0.50), "HireLens / ResumeIQ", "技能同义归一 + 分维打分")
    add_label(slide, x0, y + Inches(0.64), cw, Inches(0.18), "公开仓库的混合打分框架。LLM 略高看项目语境，确定性 40% 把幻觉拉住。", size=9, color=MUTED)
    down_arrow(slide, mid, y + h + Inches(0.0), "先过闸门，再打分")

    y = Inches(2.02)
    h = Inches(0.78)
    lane_shell(y, h, "硬门槛")
    four_w = (cw - Inches(0.28) - gap * 3) / 4
    add_chip(slide, x0, y + Inches(0.12), four_w, Inches(0.52), "年限", "低于要求则不过门")
    add_chip(slide, x0 + four_w + gap, y + Inches(0.12), four_w, Inches(0.52), "学历", "低于要求则不过门")
    add_chip(slide, x0 + (four_w + gap) * 2, y + Inches(0.12), four_w, Inches(0.52), "必备覆盖", "默认 50% 过门，缺口仍罚分")
    add_chip(slide, x0 + (four_w + gap) * 3, y + Inches(0.12), four_w, Inches(0.52), "不过门 → 不匹配", "模型抬不了红灯", accent=True)
    down_arrow(slide, mid, y + h + Inches(0.0), "过门后两路并行")

    y = Inches(3.20)
    h = Inches(2.22)
    lane_shell(y, h, "双路")
    panel_y = y + Inches(0.1)
    panel_h = Inches(2.02)
    panel_w = (cw - Inches(0.28) - Inches(0.12)) / 2
    left_x = x0
    right_x = x0 + panel_w + Inches(0.12)

    def path_panel(x, title: str, source: str) -> None:
        add_round(slide, x, panel_y, panel_w, panel_h, fill=MINT, line=INK, weight=1.5, adj=0.05)
        add_label(slide, x + Inches(0.12), panel_y + Inches(0.06), panel_w - Inches(0.24), Inches(0.22), title, size=13, bold=True, color=INK)
        add_label(slide, x + Inches(0.12), panel_y + Inches(0.26), panel_w - Inches(0.24), Inches(0.18), source, size=10, color=MUTED)

    path_panel(left_x, "LLM Judge 路", "ema-resume-ranker · 语境判断")
    path_panel(right_x, "Deterministic 路", "HireLens 分维 + ResuRank 文本")

    row_h = Inches(0.46)
    row0 = panel_y + Inches(0.50)
    inner_w = panel_w - Inches(0.24)
    add_chip(slide, left_x + Inches(0.12), row0, inner_w, row_h, "分维打分", "技能 / 经验 / 项目 / 风险")
    add_chip(slide, left_x + Inches(0.12), row0 + row_h + Inches(0.06), inner_w, row_h, "必须引用原文", "只接受 JD 或简历，网页不算")
    add_chip(slide, left_x + Inches(0.12), row0 + (row_h + Inches(0.06)) * 2, inner_w, row_h, "±18 clamp", "不能偏离确定性锚点太远")

    add_chip(slide, right_x + Inches(0.12), row0, inner_w, row_h, "技能 40%", "必备 70% + 加分 30%，先同义归一")
    add_chip(slide, right_x + Inches(0.12), row0 + row_h + Inches(0.06), inner_w, row_h, "经验 20% · 学历 10%", "可解释、可复核")
    add_chip(slide, right_x + Inches(0.12), row0 + (row_h + Inches(0.06)) * 2, inner_w, row_h, "文本 30%", "0.60 语义 + 0.40 TF-IDF")

    down_arrow(slide, mid, y + h + Inches(0.0), "加权合成，再映射决策")

    y = Inches(5.82)
    h = Inches(1.46)
    lane_shell(y, h, "合成")
    add_chip(
        slide,
        x0,
        y + Inches(0.1),
        cw - Inches(0.28),
        Inches(0.46),
        "总分 = 0.60 × 模型评委 + 0.40 × 确定性分",
        "ema-resume-ranker Ensemble",
        accent=True,
    )
    add_chip(slide, x0, y + Inches(0.64), four_w, Inches(0.48), "≥ 75 推荐", "绿灯，仍要证据撑住")
    add_chip(slide, x0 + four_w + gap, y + Inches(0.64), four_w, Inches(0.48), "60–75 复核", "黄灯，交给招聘人员")
    add_chip(slide, x0 + (four_w + gap) * 2, y + Inches(0.64), four_w, Inches(0.48), "< 60 或不匹配", "不过门槛或分不够")
    add_chip(slide, x0 + (four_w + gap) * 3, y + Inches(0.64), four_w, Inches(0.48), "Checker 只降不升", "推荐可变复核，不能加分", accent=True)
    add_label(
        slide,
        x0,
        y + Inches(1.18),
        cw,
        Inches(0.22),
        "开源给打分框架。我们加 clamp、证据绑定，以及质检失败时 fail-closed。",
        size=9,
        color=MUTED,
    )


def _cell(cell, text: str, *, size=10, bold=False, color=INK, fill: RGBColor | None = None) -> None:
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAPER
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.clear()
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    style_run(run, size=size, bold=bold, color=color)


def scoring_sources(prs: Presentation, blank) -> None:
    slide = prs.slides.add_slide(blank)
    fill_rgb(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, emu(prs.slide_width), emu(prs.slide_height)), PAPER)
    no_line(slide.shapes[-1])

    mark = add_round(slide, Inches(0.42), Inches(0.22), Inches(0.28), Inches(0.28), fill=INK, adj=0.5)
    write_shape(mark, [("简", 11, True, PAPER)], align=PP_ALIGN.CENTER, pad=0.0)
    add_label(slide, Inches(0.78), Inches(0.18), Inches(8), Inches(0.18), "简历中台 · 评测分权重依据对照", size=13, bold=True, color=INK, anchor=MSO_ANCHOR.TOP)
    add_label(
        slide,
        Inches(0.78),
        Inches(0.36),
        Inches(11.5),
        Inches(0.18),
        "主公式直接采用开源 Ensemble；分维是方法对齐；闸门 / clamp / 质检是我方 Harness。",
        size=10,
        color=MUTED,
        anchor=MSO_ANCHOR.TOP,
    )
    hair = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        emu(Inches(0.42)),
        emu(Inches(0.64)),
        emu(Inches(12.5)),
        emu(Inches(0.012)),
    )
    fill_rgb(hair, LINE)
    no_line(hair)

    rows = [
        ["评测项", "我们的权重", "来源", "对应方式"],
        ["总分 Ensemble", "0.60 LLM + 0.40 确定性", "ema-resume-ranker", "直接采用主公式"],
        ["文本分", "0.60 语义 + 0.40 TF-IDF", "ResuRank", "直接采用文本公式"],
        ["低重合文本分", "Jaccard<0.08 → 0.35 / 0.65", "ResuRank divergence", "方法对齐（未照搬压到 10%）"],
        ["技能覆盖", "必备 70% + 加分 30%", "HireLens / 常见 ATS", "方法对齐：先归一再覆盖"],
        ["确定性合成", "技能40 经验20 学历10 文本30", "ai-resume-screener 分维", "方法对齐：先合成一路再 Ensemble"],
        ["技能同义", "别名表 → 规范名", "HireLens ESCO / ResumeIQ", "方法对齐，词表自行维护"],
        ["硬门槛", "不过门 = 不匹配，不进加权", "常见 ATS Gate", "方法对齐"],
        ["决策阈值", "75 推荐 / 60 复核", "评测协议 scoring_ref", "我方标定，可配置"],
        ["LLM clamp", "相对确定性锚点 ±18", "我方 Harness", "我方约束，防模型带跑"],
        ["Checker", "recommend 只降成 review", "我方 Harness", "我方约束，不改开源公式"],
    ]
    table_shape = slide.shapes.add_table(
        len(rows),
        4,
        emu(Inches(0.42)),
        emu(Inches(0.82)),
        emu(Inches(12.5)),
        emu(Inches(6.2)),
    )
    table = table_shape.table
    table.columns[0].width = emu(Inches(2.35))
    table.columns[1].width = emu(Inches(3.55))
    table.columns[2].width = emu(Inches(3.35))
    table.columns[3].width = emu(Inches(3.25))
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                _cell(cell, text, size=11, bold=True, color=PAPER, fill=INK)
            else:
                kind = row[3]
                fill = MINT if "直接采用" in kind else (SOFT if "我方" in kind else PAPER)
                _cell(cell, text, size=11, bold=(c == 0), color=INK, fill=fill)

    add_label(
        slide,
        Inches(0.42),
        Inches(7.12),
        Inches(12.5),
        Inches(0.22),
        "不是用录用数据训出的最优权重。评测验证的是门槛、排序方向和降级规则，与 matching_eval.scoring_ref 同一把尺子。",
        size=10,
        color=MUTED,
    )


def main() -> None:
    out = Path(__file__).resolve().parents[1] / "docs" / "agent-workflow.pptx"
    prs = build()
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
